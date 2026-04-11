#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_tool.py — CLI-инструмент для создания, чтения и редактирования
PowerPoint-презентаций (.pptx) с использованием python-pptx.

Зависимости:
    pip install python-pptx Pillow
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(
        "Ошибка: библиотека python-pptx не установлена.\n"
        "Установите её командой:  pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)

try:
    from PIL import Image as PILImage  # noqa: F401
except ImportError:
    PILImage = None  # images will still work via python-pptx but without pre-validation

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

WIDESCREEN_WIDTH = Inches(13.333)
WIDESCREEN_HEIGHT = Inches(7.5)

TEMPLATES: dict[str, dict[str, Any]] = {
    "business": {
        "header_color": RGBColor(0x1A, 0x36, 0x5D),
        "accent_color": RGBColor(0x2B, 0x6C, 0xB0),
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x1A, 0x1A, 0x1A),
        "subtitle_color": RGBColor(0x4A, 0x5A, 0x6A),
        "font_name": "Calibri",
        "title_size": Pt(36),
        "subtitle_size": Pt(20),
        "body_size": Pt(18),
        "bullet_size": Pt(16),
    },
    "marketing": {
        "header_color": RGBColor(0xE5, 0x3E, 0x3E),
        "accent_color": RGBColor(0xF5, 0x65, 0x65),
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x2D, 0x2D, 0x2D),
        "subtitle_color": RGBColor(0x71, 0x71, 0x71),
        "font_name": "Calibri",
        "title_size": Pt(40),
        "subtitle_size": Pt(22),
        "body_size": Pt(18),
        "bullet_size": Pt(16),
    },
    "report": {
        "header_color": RGBColor(0x4A, 0x4A, 0x4A),
        "accent_color": RGBColor(0x9E, 0x9E, 0x9E),
        "bg_color": RGBColor(0xF7, 0xF7, 0xF7),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "subtitle_color": RGBColor(0x66, 0x66, 0x66),
        "font_name": "Calibri",
        "title_size": Pt(34),
        "subtitle_size": Pt(18),
        "body_size": Pt(16),
        "bullet_size": Pt(14),
    },
}

DEFAULT_TEMPLATE = "business"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _get_template(name: str | None) -> dict[str, Any]:
    """Return a template dict by name, falling back to DEFAULT_TEMPLATE."""
    if name is None:
        name = DEFAULT_TEMPLATE
    name = name.lower().strip()
    if name not in TEMPLATES:
        available = ", ".join(sorted(TEMPLATES))
        _die(f"Неизвестный шаблон «{name}». Доступные шаблоны: {available}")
    return TEMPLATES[name]


def _die(msg: str, code: int = 1) -> None:
    print(f"Ошибка: {msg}", file=sys.stderr)
    sys.exit(code)


def _parse_json(raw: str, label: str = "JSON") -> Any:
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        _die(f"Некорректный {label}: {exc}")


def _ensure_file_exists(path: str) -> None:
    if not os.path.isfile(path):
        _die(f"Файл не найден: {path}")


def _validate_image_path(path: str) -> str:
    """Validate that an image file exists and return its absolute path."""
    abs_path = os.path.abspath(path)
    if not os.path.isfile(abs_path):
        _die(f"Файл изображения не найден: {path}")
    return abs_path


def _set_slide_background(slide, color: RGBColor) -> None:
    """Set solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, tmpl: dict[str, Any]) -> None:
    """Add a colored header bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Emu(0),
        top=Emu(0),
        width=WIDESCREEN_WIDTH,
        height=Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = tmpl["header_color"]
    shape.line.fill.background()  # no border


def _add_title_textbox(
    slide,
    text: str,
    tmpl: dict[str, Any],
    *,
    left: float = 0.6,
    top: float = 0.15,
    width: float = 12.0,
    height: float = 0.9,
    color: RGBColor | None = None,
    size: Pt | None = None,
    bold: bool = True,
) -> None:
    """Add a title text box on the slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = tmpl["font_name"]
    p.font.size = size or tmpl["title_size"]
    p.font.bold = bold
    p.font.color.rgb = color or RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT


def _add_body_text(
    slide,
    text: str,
    tmpl: dict[str, Any],
    *,
    left: float = 0.6,
    top: float = 1.6,
    width: float = 11.5,
    height: float = 5.0,
) -> None:
    """Add body content text on the slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = tmpl["font_name"]
    p.font.size = tmpl["body_size"]
    p.font.color.rgb = tmpl["text_color"]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(8)
    return txBox


def _add_bullets(
    slide,
    bullets: list[str],
    tmpl: dict[str, Any],
    *,
    left: float = 0.6,
    top: float | None = None,
    width: float = 11.5,
    height: float = 4.5,
    start_top: float = 1.6,
) -> None:
    """Add bullet-pointed text on the slide."""
    actual_top = top if top is not None else start_top
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(actual_top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.name = tmpl["font_name"]
        p.font.size = tmpl["bullet_size"]
        p.font.color.rgb = tmpl["text_color"]
        p.level = 0
        p.space_after = Pt(6)
        # Manual bullet character
        p.text = f"\u2022  {bullet}"


def _add_image_to_slide(
    slide,
    image_path: str,
    *,
    left: float = 8.0,
    top: float = 1.8,
    max_width: float = 4.5,
    max_height: float = 4.5,
) -> None:
    """Add an image to the slide, scaling proportionally."""
    abs_path = _validate_image_path(image_path)

    # Determine original dimensions for aspect ratio
    img_width = Inches(max_width)
    img_height = Inches(max_height)

    if PILImage is not None:
        try:
            with PILImage.open(abs_path) as img:
                orig_w, orig_h = img.size
                ratio = min(max_width / orig_w, max_height / orig_h)
                img_width = Inches(orig_w * ratio)
                img_height = Inches(orig_h * ratio)
        except Exception:
            pass  # fall back to max dimensions

    slide.shapes.add_picture(
        abs_path, Inches(left), Inches(top), img_width, img_height
    )


# ---------------------------------------------------------------------------
# Title slide builder
# ---------------------------------------------------------------------------


def _build_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str | None,
    tmpl: dict[str, Any],
) -> None:
    """Build a dedicated title slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["bg_color"])

    # Full-height accent bar on the left
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Emu(0),
        top=Emu(0),
        width=Inches(0.35),
        height=WIDESCREEN_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = tmpl["header_color"]
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = tmpl["font_name"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = tmpl["header_color"]
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.2), Inches(11.0), Inches(1.2)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = tmpl["font_name"]
        p2.font.size = tmpl["subtitle_size"]
        p2.font.color.rgb = tmpl["subtitle_color"]
        p2.alignment = PP_ALIGN.LEFT

    # Decorative line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(1.0),
        top=Inches(3.9),
        width=Inches(3.0),
        height=Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = tmpl["accent_color"]
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Content slide builder
# ---------------------------------------------------------------------------


def _build_two_column_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    tmpl: dict[str, Any],
) -> None:
    """Build a two-column slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["bg_color"])
    _add_header_bar(slide, tmpl)

    slide_title = slide_data.get("title", "")
    if slide_title:
        _add_title_textbox(slide, slide_title, tmpl)

    left_text = slide_data.get("left", "")
    right_text = slide_data.get("right", "")

    if left_text:
        _add_body_text(slide, left_text, tmpl, left=0.6, top=1.6, width=5.5, height=5.0)
    if right_text:
        _add_body_text(slide, right_text, tmpl, left=6.8, top=1.6, width=5.5, height=5.0)

    # Vertical separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(6.5), top=Inches(1.4),
        width=Pt(2), height=Inches(5.2),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = tmpl["accent_color"]
    sep.line.fill.background()


def _build_section_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    tmpl: dict[str, Any],
) -> None:
    """Build a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["header_color"])

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")

    if title:
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.5), Inches(11.0), Inches(2.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = tmpl["font_name"]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.8), Inches(11.0), Inches(1.0)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = tmpl["font_name"]
        p2.font.size = tmpl["subtitle_size"]
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p2.alignment = PP_ALIGN.CENTER


def _build_quote_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    tmpl: dict[str, Any],
) -> None:
    """Build a quote slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["bg_color"])

    quote = slide_data.get("quote", slide_data.get("title", ""))
    author = slide_data.get("author", slide_data.get("subtitle", ""))

    # Large opening quote mark
    q_mark = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(2.0), Inches(2.0)
    )
    tf_q = q_mark.text_frame
    p_q = tf_q.paragraphs[0]
    p_q.text = "\u201C"
    p_q.font.size = Pt(120)
    p_q.font.color.rgb = tmpl["accent_color"]
    p_q.font.bold = True

    # Quote text
    if quote:
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(10.0), Inches(3.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        p.font.name = tmpl["font_name"]
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = tmpl["text_color"]
        p.alignment = PP_ALIGN.CENTER

    # Author
    if author:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.8), Inches(10.0), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = f"\u2014 {author}"
        p2.font.name = tmpl["font_name"]
        p2.font.size = tmpl["subtitle_size"]
        p2.font.color.rgb = tmpl["subtitle_color"]
        p2.alignment = PP_ALIGN.CENTER


def _build_table_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    tmpl: dict[str, Any],
) -> None:
    """Build a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["bg_color"])
    _add_header_bar(slide, tmpl)

    slide_title = slide_data.get("title", "")
    if slide_title:
        _add_title_textbox(slide, slide_title, tmpl)

    rows_data = slide_data.get("rows", [])
    if not rows_data or not isinstance(rows_data, list):
        return

    num_rows = len(rows_data)
    num_cols = max(len(r) for r in rows_data) if rows_data else 1

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.6), Inches(1.6),
        Inches(12.0), Inches(min(num_rows * 0.6, 5.5)),
    )
    table = table_shape.table

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            if c_idx < num_cols:
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = tmpl["font_name"]
                    paragraph.font.size = tmpl["bullet_size"]
                    paragraph.font.color.rgb = tmpl["text_color"]
                    if r_idx == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Style header row
    for c_idx in range(num_cols):
        cell = table.cell(0, c_idx)
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = tmpl["header_color"]


def _build_content_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    tmpl: dict[str, Any],
) -> None:
    """Build a single content slide from a slide data dict.
    Dispatches to specialized builders based on 'layout' field.
    """
    layout = slide_data.get("layout", "content")

    # Dispatch to specialized layout builders
    if layout == "title":
        _build_title_slide(
            prs,
            slide_data.get("title", ""),
            slide_data.get("subtitle"),
            tmpl,
        )
        return
    if layout == "two_column":
        _build_two_column_slide(prs, slide_data, tmpl)
        return
    if layout == "section":
        _build_section_slide(prs, slide_data, tmpl)
        return
    if layout == "quote":
        _build_quote_slide(prs, slide_data, tmpl)
        return
    if layout == "table":
        _build_table_slide(prs, slide_data, tmpl)
        return
    if layout == "blank":
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_background(sl, tmpl["bg_color"])
        return

    # Default: content / image / image_text / chart_placeholder
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, tmpl["bg_color"])

    slide_title = slide_data.get("title", "")
    # Accept both "content" and "body" field names for compatibility
    content = slide_data.get("content", "") or slide_data.get("body", "")
    bullets = slide_data.get("bullets", [])
    image = slide_data.get("image")
    caption = slide_data.get("caption", "")

    # Header bar
    _add_header_bar(slide, tmpl)

    # Title on the header
    if slide_title:
        _add_title_textbox(slide, slide_title, tmpl)

    # Determine layout: if there's an image, text goes left; otherwise full width
    has_image = image is not None and image != ""
    text_width = 6.5 if has_image else 11.5

    current_top = 1.6

    # Body text
    if content:
        _add_body_text(
            slide, content, tmpl, width=text_width, top=current_top
        )
        # Rough estimate of how far down the content extends
        lines_est = max(1, len(content) // 80 + 1)
        current_top += 0.5 + lines_est * 0.35

    # Bullets
    if bullets:
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        _add_bullets(
            slide,
            [str(b) for b in bullets],
            tmpl,
            width=text_width,
            start_top=current_top,
        )

    # Image
    if has_image:
        _add_image_to_slide(slide, image)

    # Caption (for image/image_text layouts)
    if caption:
        _add_body_text(
            slide, caption, tmpl,
            left=0.6, top=6.5, width=12.0, height=0.8,
        )


# ---------------------------------------------------------------------------
# Subcommand: create
# ---------------------------------------------------------------------------


def cmd_create(args: argparse.Namespace) -> None:
    tmpl = _get_template(args.template)

    prs = Presentation()
    prs.slide_width = WIDESCREEN_WIDTH
    prs.slide_height = WIDESCREEN_HEIGHT

    # Title slide
    _build_title_slide(prs, args.title or "Презентация", args.subtitle, tmpl)

    # Additional slides from --slides JSON
    if args.slides:
        slides_data = _parse_json(args.slides, "слайдов (--slides)")
        if not isinstance(slides_data, list):
            _die("--slides должен содержать JSON-массив объектов слайдов.")
        for idx, sd in enumerate(slides_data, start=1):
            if not isinstance(sd, dict):
                _die(
                    f"Элемент #{idx} в --slides должен быть объектом (dict), "
                    f"получено: {type(sd).__name__}"
                )
            _build_content_slide(prs, sd, tmpl)

    output = args.output or "presentation.pptx"
    if not output.lower().endswith(".pptx"):
        output += ".pptx"

    # Ensure parent directory exists
    parent = os.path.dirname(os.path.abspath(output))
    os.makedirs(parent, exist_ok=True)

    prs.save(output)
    print(f"Презентация сохранена: {os.path.abspath(output)}")


# ---------------------------------------------------------------------------
# Subcommand: read
# ---------------------------------------------------------------------------


def _extract_slides(prs: Presentation) -> list[dict[str, Any]]:
    """Extract text content from all slides."""
    result = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        result.append({"slide": slide_num, "texts": texts})
    return result


def cmd_read(args: argparse.Namespace) -> None:
    _ensure_file_exists(args.input)
    prs = Presentation(args.input)
    slides = _extract_slides(prs)

    for s in slides:
        print(f"--- Слайд {s['slide']} ---")
        if s["texts"]:
            for t in s["texts"]:
                print(f"  {t}")
        else:
            print("  (пусто)")
        print()


# ---------------------------------------------------------------------------
# Subcommand: edit
# ---------------------------------------------------------------------------


def cmd_edit(args: argparse.Namespace) -> None:
    _ensure_file_exists(args.input)

    prs = Presentation(args.input)
    total = len(prs.slides)
    slide_idx = args.slide - 1  # convert to 0-based

    if slide_idx < 0 or slide_idx >= total:
        _die(
            f"Номер слайда {args.slide} вне диапазона. "
            f"В презентации {total} слайд(ов)."
        )

    slide = prs.slides[slide_idx]

    # Strategy: find the first text-frame-bearing shape for title,
    # the second for body content.  This works for slides built by this tool
    # and for most standard layouts.
    text_shapes = [sh for sh in slide.shapes if sh.has_text_frame]

    if args.title is not None:
        if not text_shapes:
            _die(
                f"На слайде {args.slide} нет текстовых элементов "
                "для замены заголовка."
            )
        tf = text_shapes[0].text_frame
        for p in tf.paragraphs:
            # Preserve formatting of the first run
            if p.runs:
                p.runs[0].text = args.title
                for run in p.runs[1:]:
                    run.text = ""
            else:
                p.text = args.title
            break  # only first paragraph

    if args.content is not None:
        target = text_shapes[1] if len(text_shapes) > 1 else (
            text_shapes[0] if text_shapes else None
        )
        if target is None:
            _die(
                f"На слайде {args.slide} нет текстовых элементов "
                "для замены содержимого."
            )
        tf = target.text_frame
        # Clear existing paragraphs and set new content
        for p in tf.paragraphs:
            if p.runs:
                p.runs[0].text = ""
                for run in p.runs[1:]:
                    run.text = ""
            else:
                p.text = ""
        # Set first paragraph to new content
        first_p = tf.paragraphs[0]
        if first_p.runs:
            first_p.runs[0].text = args.content
        else:
            first_p.text = args.content

    output = args.output or args.input
    if not output.lower().endswith(".pptx"):
        output += ".pptx"

    parent = os.path.dirname(os.path.abspath(output))
    os.makedirs(parent, exist_ok=True)

    prs.save(output)
    print(f"Слайд {args.slide} отредактирован. Сохранено: {os.path.abspath(output)}")


# ---------------------------------------------------------------------------
# Subcommand: add-slide
# ---------------------------------------------------------------------------


def cmd_add_slide(args: argparse.Namespace) -> None:
    _ensure_file_exists(args.input)

    prs = Presentation(args.input)

    # Detect template from existing header color, or default
    tmpl = _get_template(args.template)

    slide_data: dict[str, Any] = {}
    if args.title:
        slide_data["title"] = args.title
    if args.content:
        slide_data["content"] = args.content
    if args.bullets:
        bullets = _parse_json(args.bullets, "буллетов (--bullets)")
        if not isinstance(bullets, list):
            _die("--bullets должен содержать JSON-массив строк.")
        slide_data["bullets"] = bullets
    if args.image:
        slide_data["image"] = args.image

    _build_content_slide(prs, slide_data, tmpl)

    output = args.output or args.input
    if not output.lower().endswith(".pptx"):
        output += ".pptx"

    parent = os.path.dirname(os.path.abspath(output))
    os.makedirs(parent, exist_ok=True)

    prs.save(output)
    total = len(prs.slides)
    print(
        f"Добавлен слайд #{total}. "
        f"Сохранено: {os.path.abspath(output)}"
    )


# ---------------------------------------------------------------------------
# Subcommand: export
# ---------------------------------------------------------------------------


def cmd_export(args: argparse.Namespace) -> None:
    _ensure_file_exists(args.input)

    prs = Presentation(args.input)
    slides = _extract_slides(prs)
    fmt = (args.format or "text").lower().strip()

    if fmt == "json":
        print(json.dumps(slides, ensure_ascii=False, indent=2))
    elif fmt == "text":
        for s in slides:
            print(f"=== Слайд {s['slide']} ===")
            if s["texts"]:
                for t in s["texts"]:
                    print(t)
            else:
                print("(пусто)")
            print()
    else:
        _die(f"Неподдерживаемый формат «{fmt}». Используйте: text, json")


# ---------------------------------------------------------------------------
# Argument parser
# ---------------------------------------------------------------------------


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="pptx_tool",
        description="CLI-инструмент для работы с PowerPoint-презентациями (.pptx)",
    )
    sub = parser.add_subparsers(dest="command", help="Доступные команды")

    # -- create --
    p_create = sub.add_parser("create", help="Создать новую презентацию")
    p_create.add_argument(
        "--output", "-o", default="presentation.pptx",
        help="Путь к выходному файлу .pptx (по умолчанию: presentation.pptx)",
    )
    p_create.add_argument(
        "--title", "-t", default="Презентация",
        help="Заголовок презентации",
    )
    p_create.add_argument(
        "--subtitle", "-s", default=None,
        help="Подзаголовок титульного слайда",
    )
    p_create.add_argument(
        "--slides", default=None,
        help=(
            'JSON-строка с массивом слайдов. Пример: '
            '\'[{"title":"Слайд 1","content":"Текст","bullets":["А","Б"]}]\''
        ),
    )
    p_create.add_argument(
        "--template", default=None,
        choices=sorted(TEMPLATES.keys()),
        help="Шаблон оформления (business, marketing, report)",
    )

    # -- read --
    p_read = sub.add_parser("read", help="Прочитать содержимое презентации")
    p_read.add_argument(
        "--input", "-i", required=True,
        help="Путь к входному файлу .pptx",
    )

    # -- edit --
    p_edit = sub.add_parser("edit", help="Редактировать слайд презентации")
    p_edit.add_argument(
        "--input", "-i", required=True,
        help="Путь к входному файлу .pptx",
    )
    p_edit.add_argument(
        "--output", "-o", default=None,
        help="Путь к выходному файлу (по умолчанию — перезапись входного)",
    )
    p_edit.add_argument(
        "--slide", "-n", type=int, required=True,
        help="Номер слайда (начиная с 1)",
    )
    p_edit.add_argument(
        "--title", "-t", default=None,
        help="Новый заголовок слайда",
    )
    p_edit.add_argument(
        "--content", "-c", default=None,
        help="Новое текстовое содержимое слайда",
    )

    # -- add-slide --
    p_add = sub.add_parser("add-slide", help="Добавить слайд в существующую презентацию")
    p_add.add_argument(
        "--input", "-i", required=True,
        help="Путь к входному файлу .pptx",
    )
    p_add.add_argument(
        "--output", "-o", default=None,
        help="Путь к выходному файлу (по умолчанию — перезапись входного)",
    )
    p_add.add_argument(
        "--title", "-t", default=None,
        help="Заголовок нового слайда",
    )
    p_add.add_argument(
        "--content", "-c", default=None,
        help="Текстовое содержимое нового слайда",
    )
    p_add.add_argument(
        "--bullets", "-b", default=None,
        help='JSON-массив буллетов. Пример: \'["Пункт 1","Пункт 2"]\'',
    )
    p_add.add_argument(
        "--image", default=None,
        help="Путь к изображению для вставки на слайд",
    )
    p_add.add_argument(
        "--template", default=None,
        choices=sorted(TEMPLATES.keys()),
        help="Шаблон оформления для нового слайда",
    )

    # -- export --
    p_export = sub.add_parser(
        "export", help="Экспортировать содержимое в текст или JSON"
    )
    p_export.add_argument(
        "--input", "-i", required=True,
        help="Путь к входному файлу .pptx",
    )
    p_export.add_argument(
        "--format", "-f", default="text", choices=["text", "json"],
        help="Формат вывода: text или json (по умолчанию: text)",
    )

    return parser


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    if args.command is None:
        parser.print_help()
        sys.exit(0)

    dispatch = {
        "create": cmd_create,
        "read": cmd_read,
        "edit": cmd_edit,
        "add-slide": cmd_add_slide,
        "export": cmd_export,
    }

    handler = dispatch.get(args.command)
    if handler is None:
        parser.print_help()
        sys.exit(1)

    try:
        handler(args)
    except KeyboardInterrupt:
        print("\nПрервано пользователем.", file=sys.stderr)
        sys.exit(130)
    except Exception as exc:
        _die(f"Непредвиденная ошибка: {exc}")


if __name__ == "__main__":
    main()
